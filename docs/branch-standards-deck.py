#!/usr/bin/env python3
"""Generate SDLC Branch & Traceability Standards PowerPoint deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- Color palette --
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x8C, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xC9, 0x7B)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA8, 0x5C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x26, 0x26, 0x3E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=18, color=WHITE, bold=False,
                  space_before=Pt(6), space_after=Pt(2), alignment=PP_ALIGN.LEFT,
                  font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    return p


def add_card(slide, left, top, width, height, color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_accent_bar(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num} / {total}", font_size=11, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 15

# ============================================================
# SLIDE 1 — Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "SDLC Branch & Traceability Standards", font_size=40, color=WHITE, bold=True)

add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
             "Branch naming conventions, PR standards, and audit traceability",
             font_size=22, color=LIGHT_GRAY)

add_accent_bar(slide, Inches(1.5), Inches(3.9), Inches(2), Inches(0.05), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(4.5), Inches(6), Inches(0.5),
             "Engineering Standards  |  Version 1.0  |  2026", font_size=14, color=MID_GRAY)

add_slide_number(slide, 1, TOTAL_SLIDES)

# ============================================================
# SLIDE 2 — Why Traceability Matters
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Why Traceability Matters", font_size=32, color=WHITE, bold=True)

cards = [
    ("Audit Compliance", "Every code change links back to a requirement,\nuser story, or defect — providing a clear audit trail\nfrom request to deployment.",
     ACCENT_BLUE),
    ("Change Impact Analysis", "Trace which features, stories, and bugs are\nincluded in any release. Quickly identify what\nchanged and why.",
     ACCENT_GREEN),
    ("Accountability", "Clear ownership of changes. Every branch, commit,\nand PR is tied to a trackable work item with an\nassigned developer.",
     ACCENT_ORANGE),
    ("Regulatory Requirements", "SOC 2, ISO 27001, HIPAA, and SOX audits\nrequire demonstrable change management with\nfull traceability.",
     ACCENT_PURPLE),
]

for i, (title, desc, accent) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.1)
    y = Inches(1.5) + row * Inches(2.7)
    add_card(slide, x, y, Inches(5.7), Inches(2.3))
    add_accent_bar(slide, x, y, Inches(0.08), Inches(2.3), accent)
    add_text_box(slide, x + Inches(0.4), y + Inches(0.3), Inches(5), Inches(0.5),
                 title, font_size=20, color=accent, bold=True)
    add_text_box(slide, x + Inches(0.4), y + Inches(0.9), Inches(5), Inches(1.2),
                 desc, font_size=14, color=LIGHT_GRAY)

add_slide_number(slide, 2, TOTAL_SLIDES)

# ============================================================
# SLIDE 3 — Protection Levels Defined
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Branch Protection Levels Defined", font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Two tiers of protection applied via GitHub branch protection rules", font_size=14, color=MID_GRAY)

# -- Fully Protected card --
add_card(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(2.5))
add_accent_bar(slide, Inches(0.8), Inches(1.6), Inches(0.08), Inches(2.5), ACCENT_RED)

shield_shape = add_card(slide, Inches(1.2), Inches(1.75), Inches(1.8), Inches(0.45), ACCENT_RED)
add_text_box(slide, Inches(1.2), Inches(1.77), Inches(1.8), Inches(0.4),
             "FULLY PROTECTED", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

prot_rules = [
    "PR required for all changes — no direct pushes or commits",
    "Minimum 1 reviewer approval before merge",
    "All CI/CD status checks must pass",
    "Force pushes and branch deletion are disabled",
    "Signed commits recommended (optional enforcement)",
]
for i, rule in enumerate(prot_rules):
    add_text_box(slide, Inches(1.5), Inches(2.35) + i * Inches(0.33), Inches(10), Inches(0.3),
                 f"\u2713  {rule}", font_size=13, color=LIGHT_GRAY)

# -- Standard Protected card --
add_card(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(2.8))
add_accent_bar(slide, Inches(0.8), Inches(4.4), Inches(0.08), Inches(2.8), ACCENT_ORANGE)

shield2 = add_card(slide, Inches(1.2), Inches(4.55), Inches(2.2), Inches(0.45), ACCENT_ORANGE)
add_text_box(slide, Inches(1.2), Inches(4.57), Inches(2.2), Inches(0.4),
             "STANDARD PROTECTED", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

std_rules = [
    "PR required for all changes — no direct pushes",
    "Minimum 1 reviewer approval before merge",
    "CI/CD status checks must pass",
    "Force pushes disabled",
    "Branch deletion IS allowed (these are temporary branches)",
    "Relaxed: may allow maintainer override for time-sensitive releases",
]
for i, rule in enumerate(std_rules):
    add_text_box(slide, Inches(1.5), Inches(5.15) + i * Inches(0.33), Inches(10), Inches(0.3),
                 f"\u2713  {rule}", font_size=13, color=LIGHT_GRAY)

# Key difference callout
add_card(slide, Inches(6.5), Inches(1.75), Inches(5.5), Inches(0.45), CARD_BG)
add_text_box(slide, Inches(6.7), Inches(1.77), Inches(5.2), Inches(0.4),
             "Applies to:  main, develop", font_size=13, color=ACCENT_RED, bold=True)

add_card(slide, Inches(6.5), Inches(4.55), Inches(5.5), Inches(0.45), CARD_BG)
add_text_box(slide, Inches(6.7), Inches(4.57), Inches(5.2), Inches(0.4),
             "Applies to:  release/*, hotfix/*", font_size=13, color=ACCENT_ORANGE, bold=True)

add_slide_number(slide, 3, TOTAL_SLIDES)

# ============================================================
# SLIDE 4 — Mandatory vs Optional Branches
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Mandatory vs Optional Branches", font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Not all long-lived branches are required — some depend on your team's Git workflow", font_size=14, color=MID_GRAY)

# -- MANDATORY section --
add_card(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.5))
add_accent_bar(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.06), ACCENT_RED)

req_badge = add_card(slide, Inches(1.2), Inches(1.8), Inches(1.8), Inches(0.45), ACCENT_RED)
add_text_box(slide, Inches(1.2), Inches(1.82), Inches(1.8), Inches(0.4),
             "REQUIRED", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# main
add_text_box(slide, Inches(1.2), Inches(2.5), Inches(4.5), Inches(0.5),
             "main", font_size=22, color=ACCENT_RED, bold=True, font_name="Consolas")
add_text_box(slide, Inches(1.2), Inches(2.95), Inches(4.8), Inches(0.4),
             "Production-ready code. Always exists.", font_size=13, color=LIGHT_GRAY)
add_text_box(slide, Inches(1.2), Inches(3.3), Inches(4.8), Inches(0.4),
             "Protection: Fully Protected", font_size=13, color=ACCENT_RED, bold=True)

# feature branches
add_accent_bar(slide, Inches(1.2), Inches(3.8), Inches(4.5), Inches(0.02), CARD_BG)
add_text_box(slide, Inches(1.2), Inches(4.0), Inches(4.5), Inches(0.5),
             "feature/* | bugfix/* | hotfix/*", font_size=18, color=ACCENT_BLUE, bold=True, font_name="Consolas")
add_text_box(slide, Inches(1.2), Inches(4.5), Inches(4.8), Inches(0.8),
             "Short-lived branches for individual work items.\nAlways branch from main or develop.\nAlways deleted after merge.", font_size=13, color=LIGHT_GRAY)
add_text_box(slide, Inches(1.2), Inches(5.3), Inches(4.8), Inches(0.4),
             "Protection: Not protected (ephemeral)", font_size=13, color=MID_GRAY, bold=True)

# Naming requirement
add_card(slide, Inches(1.0), Inches(5.8), Inches(5.0), Inches(1.0), CARD_BG)
add_text_box(slide, Inches(1.2), Inches(5.85), Inches(4.6), Inches(0.9),
             "Ticket ID in branch name is MANDATORY\nfor all feature/bugfix/hotfix branches.\nThis is the primary traceability link.",
             font_size=12, color=ACCENT_ORANGE, bold=True)

# -- OPTIONAL section --
add_card(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(5.5))
add_accent_bar(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.06), ACCENT_GREEN)

opt_badge = add_card(slide, Inches(7.2), Inches(1.8), Inches(2.5), Inches(0.45), ACCENT_GREEN)
add_text_box(slide, Inches(7.2), Inches(1.82), Inches(2.5), Inches(0.4),
             "OPTIONAL (if used, protect it)", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# develop
add_text_box(slide, Inches(7.2), Inches(2.5), Inches(4.5), Inches(0.5),
             "develop", font_size=22, color=ACCENT_ORANGE, bold=True, font_name="Consolas")
add_text_box(slide, Inches(7.2), Inches(2.95), Inches(4.8), Inches(1.0),
             "Integration branch for teams using Gitflow.\nFeature branches merge here before release.\nIf your team uses trunk-based development\n(feature branches merge directly to main),\nyou do not need this branch.", font_size=13, color=LIGHT_GRAY)
add_text_box(slide, Inches(7.2), Inches(3.95), Inches(4.8), Inches(0.4),
             "If it exists: Fully Protected", font_size=13, color=ACCENT_ORANGE, bold=True)

# release
add_accent_bar(slide, Inches(7.2), Inches(4.5), Inches(4.5), Inches(0.02), CARD_BG)
add_text_box(slide, Inches(7.2), Inches(4.7), Inches(4.5), Inches(0.5),
             "release/*", font_size=22, color=ACCENT_GREEN, bold=True, font_name="Consolas")
add_text_box(slide, Inches(7.2), Inches(5.15), Inches(4.8), Inches(1.0),
             "Stabilization branch created before a release.\nOnly needed if you separate release prep from\nnormal development. Temporary — deleted\nafter merging to main.", font_size=13, color=LIGHT_GRAY)
add_text_box(slide, Inches(7.2), Inches(6.15), Inches(4.8), Inches(0.4),
             "If it exists: Standard Protected", font_size=13, color=ACCENT_GREEN, bold=True)

# Bottom callout
add_card(slide, Inches(7.0), Inches(6.6), Inches(5.2), Inches(0.45), CARD_BG)
add_text_box(slide, Inches(7.2), Inches(6.62), Inches(4.8), Inches(0.4),
             "Rule: If a branch exists, it must be protected.", font_size=13, color=ACCENT_RED, bold=True)

add_slide_number(slide, 4, TOTAL_SLIDES)

# ============================================================
# SLIDE 5 — Branch Naming Convention
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Branch Naming Convention", font_size=32, color=WHITE, bold=True)

# Formula
add_card(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.5))
add_text_box(slide, Inches(1.2), Inches(1.35), Inches(10.5), Inches(0.4),
             "Standard Pattern", font_size=13, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(1.2), Inches(1.7), Inches(10.5), Inches(0.5),
             "<type>/<ticket-id>-<short-description>", font_size=24, color=ACCENT_GREEN,
             bold=True, font_name="Consolas")

add_text_box(slide, Inches(1.2), Inches(2.15), Inches(10.5), Inches(0.4),
             "Optional — with PI and Sprint", font_size=13, color=ACCENT_PURPLE, bold=True)
add_text_box(slide, Inches(1.2), Inches(2.45), Inches(10.5), Inches(0.4),
             "<type>/pi<N>-s<N>/<ticket-id>-<short-description>", font_size=20, color=ACCENT_PURPLE,
             bold=True, font_name="Consolas")

# Types table
types = [
    ("feature/", "New functionality tied to a user story", "feature/IS-1234-add-user-profile", ACCENT_BLUE),
    ("integration/", "Multi-story feature (integration branch)", "integration/IS-200-user-onboarding", ACCENT_PURPLE),
    ("bugfix/", "Defect fix from backlog", "bugfix/IS-567-fix-login-timeout", ACCENT_ORANGE),
    ("hotfix/", "Urgent production fix (branches from main)", "hotfix/IS-890-patch-auth-bypass", ACCENT_RED),
    ("chore/", "Tech debt, refactoring, CI/CD updates", "chore/IS-321-upgrade-spring-boot", MID_GRAY),
    ("release/", "Release stabilization branch", "release/v2.4.0", ACCENT_GREEN),
]

y_start = Inches(3.1)
# Header
add_text_box(slide, Inches(0.8), y_start, Inches(2), Inches(0.4),
             "PREFIX", font_size=11, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(3.0), y_start, Inches(3.5), Inches(0.4),
             "USE CASE", font_size=11, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(6.8), y_start, Inches(6), Inches(0.4),
             "EXAMPLE  (with optional pi/sprint)", font_size=11, color=ACCENT_BLUE, bold=True)

add_accent_bar(slide, Inches(0.8), y_start + Inches(0.35), Inches(11.5), Inches(0.02), CARD_BG)

for i, (prefix, use, example, color) in enumerate(types):
    y = y_start + Inches(0.5) + i * Inches(0.55)
    if i % 2 == 0:
        add_card(slide, Inches(0.8), y - Inches(0.05), Inches(11.5), Inches(0.52), CARD_BG)
    add_text_box(slide, Inches(0.8), y, Inches(2), Inches(0.5),
                 prefix, font_size=14, color=color, bold=True, font_name="Consolas")
    add_text_box(slide, Inches(3.0), y + Inches(0.03), Inches(3.5), Inches(0.5),
                 use, font_size=12, color=LIGHT_GRAY)
    add_text_box(slide, Inches(6.8), y + Inches(0.03), Inches(6), Inches(0.5),
                 example, font_size=12, color=WHITE, font_name="Consolas")

# PI/Sprint example row
y_pi = y_start + Inches(0.5) + len(types) * Inches(0.55) + Inches(0.15)
add_card(slide, Inches(0.8), y_pi - Inches(0.05), Inches(11.5), Inches(0.52), CARD_BG)
add_text_box(slide, Inches(0.8), y_pi, Inches(2), Inches(0.5),
             "with PI/Sprint", font_size=12, color=ACCENT_PURPLE, bold=True)
add_text_box(slide, Inches(3.0), y_pi + Inches(0.03), Inches(3.5), Inches(0.5),
             "Optional namespace for SAFe teams", font_size=12, color=LIGHT_GRAY)
add_text_box(slide, Inches(6.8), y_pi + Inches(0.03), Inches(6), Inches(0.5),
             "feature/pi3-s2/IS-1234-add-user-profile", font_size=12, color=ACCENT_PURPLE, font_name="Consolas")

# Rules
add_card(slide, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.55))
tf = add_text_box(slide, Inches(1.2), Inches(6.72), Inches(10.5), Inches(0.5),
                  "Rules: ", font_size=13, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "Lowercase only  |  Hyphens as separators  |  IS-xxxx ticket ID mandatory  |  PI/Sprint optional  |  Delete after merge",
              font_size=12, color=LIGHT_GRAY)

add_slide_number(slide, 5, TOTAL_SLIDES)

# ============================================================
# SLIDE 6 — PR Naming Convention
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Pull Request Standards", font_size=32, color=WHITE, bold=True)

# PR Title Formula
add_card(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.0))
add_text_box(slide, Inches(1.2), Inches(1.35), Inches(10.5), Inches(0.4),
             "PR Title Pattern", font_size=13, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(1.2), Inches(1.7), Inches(10.5), Inches(0.5),
             "<type>(<scope>): <ticket-id> <description>", font_size=24, color=ACCENT_GREEN,
             bold=True, font_name="Consolas")

# Examples
examples = [
    ("feat(auth): IS-1234 Add OAuth2 login flow", ACCENT_BLUE),
    ("fix(api): IS-567 Resolve timeout on bulk export", ACCENT_ORANGE),
    ("hotfix(payments): IS-890 Patch null pointer in refund handler", ACCENT_RED),
    ("chore(ci): IS-321 Upgrade to Gradle 8.12", MID_GRAY),
]

y_start = Inches(2.6)
add_text_box(slide, Inches(0.8), y_start, Inches(3), Inches(0.4),
             "EXAMPLES", font_size=11, color=ACCENT_BLUE, bold=True)

for i, (ex, color) in enumerate(examples):
    y = y_start + Inches(0.45) + i * Inches(0.6)
    add_accent_bar(slide, Inches(0.8), y + Inches(0.05), Inches(0.06), Inches(0.35), color)
    add_text_box(slide, Inches(1.1), y, Inches(10.5), Inches(0.5),
                 ex, font_size=15, color=WHITE, font_name="Consolas")

# PR Requirements
y_req = Inches(5.0)
add_text_box(slide, Inches(0.8), y_req, Inches(6), Inches(0.4),
             "PR REQUIREMENTS", font_size=11, color=ACCENT_BLUE, bold=True)

reqs_left = [
    "Link to ticket/user story in description",
    "At least 1 reviewer approval required",
    "All CI checks must pass",
    "No merge conflicts with target branch",
]
reqs_right = [
    "Squash-merge to keep history clean",
    "Delete source branch after merge",
    "Include test evidence or test plan",
    "Scope: one ticket per PR (preferred)",
]

for i, req in enumerate(reqs_left):
    y = y_req + Inches(0.4) + i * Inches(0.45)
    add_text_box(slide, Inches(1.1), y, Inches(5.5), Inches(0.4),
                 f"\u2713  {req}", font_size=13, color=LIGHT_GRAY)

for i, req in enumerate(reqs_right):
    y = y_req + Inches(0.4) + i * Inches(0.45)
    add_text_box(slide, Inches(7.0), y, Inches(5.5), Inches(0.4),
                 f"\u2713  {req}", font_size=13, color=LIGHT_GRAY)

add_slide_number(slide, 6, TOTAL_SLIDES)

# ============================================================
# SLIDE 7 — Commit Message Standard
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Commit Message Standard", font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Follow Conventional Commits (conventionalcommits.org)", font_size=14, color=MID_GRAY)

# Format
add_card(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.8))
add_text_box(slide, Inches(1.2), Inches(1.6), Inches(10), Inches(0.4),
             "Format", font_size=13, color=ACCENT_BLUE, bold=True)

lines = [
    ("<type>(<scope>): <short summary>", ACCENT_GREEN),
    ("", WHITE),
    ("[optional body — explain WHY, not WHAT]", LIGHT_GRAY),
    ("", WHITE),
    ("[optional footer]", LIGHT_GRAY),
    ("Refs: IS-1234", ACCENT_BLUE),
    ("BREAKING CHANGE: <description>", ACCENT_RED),
]

for i, (line, color) in enumerate(lines):
    add_text_box(slide, Inches(1.5), Inches(2.0) + i * Inches(0.35), Inches(9), Inches(0.35),
                 line, font_size=15, color=color, font_name="Consolas")

# Types — left column
y_types = Inches(4.5)
add_text_box(slide, Inches(0.8), y_types, Inches(4), Inches(0.4),
             "ALLOWED TYPES", font_size=11, color=ACCENT_BLUE, bold=True)

type_items = [
    ("feat", "New feature"), ("fix", "Bug fix"), ("docs", "Documentation"),
    ("style", "Formatting"), ("refactor", "Restructure"), ("test", "Add/fix tests"),
    ("chore", "Maintenance"), ("perf", "Performance"), ("ci", "CI/CD changes"),
]

for i, (t, d) in enumerate(type_items):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(2.1)
    y = y_types + Inches(0.35) + row * Inches(0.38)
    add_text_box(slide, x, y, Inches(1.0), Inches(0.35),
                 t, font_size=12, color=ACCENT_GREEN, bold=True, font_name="Consolas")
    add_text_box(slide, x + Inches(1.0), y, Inches(1.0), Inches(0.35),
                 f"- {d}", font_size=11, color=LIGHT_GRAY)

# Automation box — right side
add_card(slide, Inches(7.0), Inches(4.5), Inches(5.5), Inches(2.7))
add_accent_bar(slide, Inches(7.0), Inches(4.5), Inches(0.08), Inches(2.7), ACCENT_PURPLE)

add_text_box(slide, Inches(7.4), Inches(4.6), Inches(4.8), Inches(0.4),
             "AUTOMATED TICKET INJECTION", font_size=13, color=ACCENT_PURPLE, bold=True)

add_text_box(slide, Inches(7.4), Inches(5.0), Inches(4.8), Inches(0.9),
             "Ticket ID in commits is OPTIONAL for developers.\nA Git hook (prepare-commit-msg) or CI step\nautomatically extracts the ticket ID from the\nbranch name and appends it to each commit.",
             font_size=13, color=LIGHT_GRAY)

add_text_box(slide, Inches(7.4), Inches(5.85), Inches(4.8), Inches(0.4),
             "How it works:", font_size=12, color=ACCENT_BLUE, bold=True)

auto_steps = [
    "1. Dev creates branch: feature/IS-1234-add-profile",
    "2. Dev commits: feat(user): add profile page",
    "3. Hook detects IS-1234 from branch name",
    "4. Commit becomes: feat(user): add profile page",
    "   Refs: IS-1234   (auto-appended)",
]
for i, step in enumerate(auto_steps):
    color = ACCENT_GREEN if "auto-appended" in step else LIGHT_GRAY
    add_text_box(slide, Inches(7.4), Inches(6.15) + i * Inches(0.28), Inches(4.8), Inches(0.28),
                 step, font_size=11, color=color, font_name="Consolas")

# Golden rule — bottom left
add_card(slide, Inches(0.8), Inches(6.5), Inches(5.8), Inches(0.7), CARD_BG)
add_text_box(slide, Inches(1.2), Inches(6.55), Inches(5.2), Inches(0.6),
             "Ticket ID in commits: Optional (manual)\n"
             "Ticket ID in branch name & PR title: Required",
             font_size=13, color=ACCENT_ORANGE, bold=True)

add_slide_number(slide, 7, TOTAL_SLIDES)

# ============================================================
# SLIDE 8 — Feature Branch Workflow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Feature Branch Workflow", font_size=32, color=WHITE, bold=True)

steps = [
    ("1", "Pick Up Story", "Assign user story in\nproject board.\nStatus: In Progress", ACCENT_BLUE),
    ("2", "Create Branch", "Branch from develop:\nfeature/IS-1234-\nadd-user-profile", ACCENT_BLUE),
    ("3", "Develop & Commit", "Commit early and often.\nConventional commits.\nRef ticket in footer.", ACCENT_GREEN),
    ("4", "Open PR", "PR to develop.\nLink story in description.\nRequest reviewer.", ACCENT_GREEN),
    ("5", "Review & Approve", "1+ approval required.\nCI checks pass.\nNo merge conflicts.", ACCENT_ORANGE),
    ("6", "Merge & Clean Up", "Squash-merge to develop.\nDelete feature branch.\nStory \u2192 Done.", ACCENT_ORANGE),
]

for i, (num, title, desc, accent) in enumerate(steps):
    x = Inches(0.5) + i * Inches(2.1)
    y = Inches(1.5)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    add_text_box(slide, x + Inches(0.65), y + Inches(0.08), Inches(0.6), Inches(0.5),
                 num, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Arrow (except last)
    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(1.4), y + Inches(0.1), Inches(0.5), Inches(0.5),
                     "\u2192", font_size=20, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Card
    add_card(slide, x, y + Inches(0.8), Inches(1.9), Inches(2.5))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.95), Inches(1.6), Inches(0.5),
                 title, font_size=15, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.4), Inches(1.6), Inches(1.5),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Hotfix callout
add_card(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(2.0))
add_text_box(slide, Inches(1.2), Inches(5.3), Inches(3), Inches(0.4),
             "HOTFIX EXCEPTION", font_size=13, color=ACCENT_RED, bold=True)

hotfix_lines = [
    "1. Branch from main (not develop):  hotfix/IS-890-patch-auth-bypass",
    "2. Fix, test, open PR to main with expedited review",
    "3. After merging to main, immediately merge main back to develop",
    "4. Hotfixes still require PR approval — no direct pushes to main",
]
for i, line in enumerate(hotfix_lines):
    add_text_box(slide, Inches(1.5), Inches(5.7) + i * Inches(0.38), Inches(10), Inches(0.35),
                 line, font_size=13, color=LIGHT_GRAY, font_name="Consolas" if "hotfix/" in line else "Calibri")

add_slide_number(slide, 8, TOTAL_SLIDES)

# ============================================================
# SLIDE 9 — Work Item Hierarchy & Branching Levels
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Work Item Hierarchy & Branch Mapping", font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Every feature has at least one user story — the ticket prefix in the branch name determines the pattern", font_size=14, color=MID_GRAY)

# -- Hierarchy diagram (left side) --
add_card(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.5))
add_accent_bar(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.06), ACCENT_PURPLE)

add_text_box(slide, Inches(1.2), Inches(1.8), Inches(4.8), Inches(0.4),
             "YOUR WORK ITEM HIERARCHY", font_size=14, color=ACCENT_PURPLE, bold=True)

hierarchy = [
    ("IS-100", "Epic", "Strategic initiative or large body of work.\nPlanning container only — never branched.",
     ACCENT_PURPLE, "No Branch", "Tracked via PR descriptions"),
    ("IS-200", "Feature", "Deliverable capability within an epic.\nAlways has 1 or more user stories.",
     ACCENT_BLUE, "Conditional", "integration/ branch if 2+ stories"),
    ("IS-1234", "User Story", "Single piece of user-facing value.\nThe primary unit of work for developers.",
     ACCENT_GREEN, "Always Branch", "feature/ branch (default level)"),
    ("IS-5678", "Task", "Technical sub-unit of a story.\nUsually tracked within the story's branch.",
     MID_GRAY, "Rarely Branch", "Only if large enough for own PR"),
]

for i, (ticket, label, desc, accent, branch_type, note) in enumerate(hierarchy):
    y = Inches(2.3) + i * Inches(1.2)
    indent = Inches(0.15) * i

    # Level card
    add_card(slide, Inches(1.0) + indent, y, Inches(5.0) - indent, Inches(1.0), CARD_BG)
    add_accent_bar(slide, Inches(1.0) + indent, y, Inches(0.06), Inches(1.0), accent)

    add_text_box(slide, Inches(1.2) + indent, y + Inches(0.05), Inches(1.5), Inches(0.35),
                 ticket, font_size=14, color=accent, bold=True, font_name="Consolas")
    add_text_box(slide, Inches(2.8) + indent, y + Inches(0.05), Inches(1.0), Inches(0.35),
                 label, font_size=13, color=accent, bold=True)
    add_text_box(slide, Inches(1.2) + indent, y + Inches(0.38), Inches(3.5), Inches(0.55),
                 desc, font_size=10, color=LIGHT_GRAY)

    # Branch type badge
    badge_color = accent if branch_type != "No Branch" else CARD_BG
    branch_badge = add_card(slide, Inches(4.2), y + Inches(0.05), Inches(1.6), Inches(0.35), badge_color)
    add_text_box(slide, Inches(4.2), y + Inches(0.07), Inches(1.6), Inches(0.3),
                 branch_type, font_size=9, color=WHITE if badge_color != CARD_BG else MID_GRAY,
                 bold=True, alignment=PP_ALIGN.CENTER)

# -- Two Patterns (right side) --
add_card(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(5.5))
add_accent_bar(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(7.2), Inches(1.8), Inches(4.8), Inches(0.4),
             "TWO BRANCHING PATTERNS (both allowed)", font_size=14, color=ACCENT_BLUE, bold=True)

add_text_box(slide, Inches(7.2), Inches(2.2), Inches(4.8), Inches(0.35),
             "The branch type prefix determines the pattern:", font_size=12, color=LIGHT_GRAY)

# Pattern A summary
add_card(slide, Inches(7.0), Inches(2.7), Inches(5.1), Inches(1.6), CARD_BG)
add_accent_bar(slide, Inches(7.0), Inches(2.7), Inches(0.06), Inches(1.6), ACCENT_GREEN)

pa_badge = add_card(slide, Inches(7.2), Inches(2.8), Inches(1.4), Inches(0.35), ACCENT_GREEN)
add_text_box(slide, Inches(7.2), Inches(2.82), Inches(1.4), Inches(0.3),
             "feature/", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Consolas")
add_text_box(slide, Inches(8.8), Inches(2.82), Inches(3.2), Inches(0.3),
             "Pattern A — Single Story", font_size=13, color=ACCENT_GREEN, bold=True)

add_text_box(slide, Inches(7.3), Inches(3.25), Inches(4.6), Inches(0.9),
             "Feature has 1 user story.\nBranch at story level, PR to main/develop.\n\n"
             "feature/IS-1234-signup-form \u2192 main",
             font_size=11, color=LIGHT_GRAY)

# Pattern B summary
add_card(slide, Inches(7.0), Inches(4.5), Inches(5.1), Inches(2.2), CARD_BG)
add_accent_bar(slide, Inches(7.0), Inches(4.5), Inches(0.06), Inches(2.2), ACCENT_BLUE)

pb_badge = add_card(slide, Inches(7.2), Inches(4.6), Inches(1.8), Inches(0.35), ACCENT_BLUE)
add_text_box(slide, Inches(7.2), Inches(4.62), Inches(1.8), Inches(0.3),
             "integration/", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Consolas")
add_text_box(slide, Inches(9.2), Inches(4.62), Inches(2.8), Inches(0.3),
             "Pattern B — Multi-Story", font_size=13, color=ACCENT_BLUE, bold=True)

add_text_box(slide, Inches(7.3), Inches(5.05), Inches(4.6), Inches(1.5),
             "Feature has 2+ user stories.\nIntegration branch at feature level.\nStory branches merge into it.\n\n"
             "integration/IS-200-onboarding \u2192 main\n"
             "  \u2190 feature/IS-1234-signup\n"
             "  \u2190 feature/IS-1235-verify",
             font_size=11, color=LIGHT_GRAY)

# Decision rule at bottom
add_card(slide, Inches(0.8), Inches(7.3), Inches(11.5), Inches(0.0), CARD_BG)  # spacer

add_slide_number(slide, 9, TOTAL_SLIDES)

# ============================================================
# SLIDE 10 — Pattern A Deep Dive: Single-Story Branching
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_GREEN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Pattern A: Single-Story Branching", font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Feature has one user story — branch directly at story level (most common pattern)", font_size=14, color=MID_GRAY)

# -- Visual diagram --
add_card(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.5))

# main/develop line
add_accent_bar(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(0.04), ACCENT_RED)
add_text_box(slide, Inches(1.5), Inches(1.85), Inches(2), Inches(0.4),
             "main / develop", font_size=14, color=ACCENT_RED, bold=True, font_name="Consolas")

# Feature branch line
add_accent_bar(slide, Inches(3.5), Inches(2.34), Inches(0.04), Inches(1.0), ACCENT_GREEN)
add_accent_bar(slide, Inches(3.5), Inches(3.3), Inches(5.5), Inches(0.04), ACCENT_GREEN)
add_accent_bar(slide, Inches(9.0), Inches(2.34), Inches(0.04), Inches(1.0), ACCENT_GREEN)

add_text_box(slide, Inches(4.0), Inches(3.0), Inches(4.5), Inches(0.4),
             "feature/IS-1234-signup-form", font_size=13, color=ACCENT_GREEN, bold=True, font_name="Consolas")

# Merge arrows
add_text_box(slide, Inches(3.2), Inches(2.55), Inches(0.5), Inches(0.4),
             "\u2193", font_size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.2), Inches(1.65), Inches(1.0), Inches(0.35),
             "branch", font_size=10, color=MID_GRAY)
add_text_box(slide, Inches(8.7), Inches(2.55), Inches(0.5), Inches(0.4),
             "\u2191", font_size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.5), Inches(1.65), Inches(1.5), Inches(0.35),
             "squash-merge PR", font_size=10, color=MID_GRAY)

# -- When to use --
add_card(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(3.0))
add_accent_bar(slide, Inches(0.8), Inches(4.3), Inches(0.08), Inches(3.0), ACCENT_GREEN)

add_text_box(slide, Inches(1.2), Inches(4.4), Inches(4.8), Inches(0.4),
             "WHEN TO USE", font_size=14, color=ACCENT_GREEN, bold=True)

when_a = [
    "Feature maps to exactly 1 user story",
    "One developer (or pair) owns the full scope",
    "Work can be completed in a single PR",
    "Story is independently releasable",
]
for i, item in enumerate(when_a):
    add_text_box(slide, Inches(1.5), Inches(4.9) + i * Inches(0.38), Inches(4.5), Inches(0.35),
                 f"\u2713  {item}", font_size=13, color=LIGHT_GRAY)

add_text_box(slide, Inches(1.2), Inches(6.5), Inches(4.8), Inches(0.4),
             "This is the default. Use this unless the feature\nhas multiple stories.", font_size=12, color=ACCENT_ORANGE, bold=True)

# -- Lifecycle --
add_card(slide, Inches(6.8), Inches(4.3), Inches(5.5), Inches(3.0))
add_accent_bar(slide, Inches(6.8), Inches(4.3), Inches(0.08), Inches(3.0), ACCENT_BLUE)

add_text_box(slide, Inches(7.2), Inches(4.4), Inches(4.8), Inches(0.4),
             "LIFECYCLE", font_size=14, color=ACCENT_BLUE, bold=True)

lifecycle_a = [
    ("1. Pick up story", "Assign IS-1234 on project board"),
    ("2. Create branch", "feature/IS-1234-signup-form from main/develop"),
    ("3. Develop", "Commit with conventional messages"),
    ("4. Open PR", "PR title: feat(auth): IS-1234 Add signup form"),
    ("5. Review + merge", "Squash-merge, delete branch, story \u2192 Done"),
]
for i, (step, detail) in enumerate(lifecycle_a):
    y = Inches(4.85) + i * Inches(0.42)
    add_text_box(slide, Inches(7.2), y, Inches(2.2), Inches(0.35),
                 step, font_size=12, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(9.4), y, Inches(2.7), Inches(0.35),
                 detail, font_size=11, color=LIGHT_GRAY)

# Traceability callout
add_card(slide, Inches(7.0), Inches(7.0), Inches(5.1), Inches(0.0), CARD_BG)

add_slide_number(slide, 10, TOTAL_SLIDES)

# ============================================================
# SLIDE 11 — Pattern B Deep Dive: Multi-Story Feature Branching
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Pattern B: Multi-Story Feature Branching", font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Feature has 2+ user stories — create a feature integration branch, story branches merge into it", font_size=14, color=MID_GRAY)

# -- Visual diagram --
add_card(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(3.2))

# main/develop line
add_accent_bar(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(0.04), ACCENT_RED)
add_text_box(slide, Inches(1.5), Inches(1.6), Inches(2), Inches(0.4),
             "main / develop", font_size=14, color=ACCENT_RED, bold=True, font_name="Consolas")

# Feature integration branch
add_accent_bar(slide, Inches(3.0), Inches(2.04), Inches(0.04), Inches(0.8), ACCENT_BLUE)
add_accent_bar(slide, Inches(3.0), Inches(2.8), Inches(7.5), Inches(0.04), ACCENT_BLUE)
add_accent_bar(slide, Inches(10.5), Inches(2.04), Inches(0.04), Inches(0.8), ACCENT_BLUE)
add_text_box(slide, Inches(4.5), Inches(2.45), Inches(5), Inches(0.4),
             "integration/IS-200-user-onboarding", font_size=13, color=ACCENT_BLUE, bold=True, font_name="Consolas")

# Story branches
add_accent_bar(slide, Inches(4.0), Inches(2.84), Inches(0.04), Inches(0.7), ACCENT_GREEN)
add_accent_bar(slide, Inches(4.0), Inches(3.5), Inches(2.5), Inches(0.04), ACCENT_GREEN)
add_accent_bar(slide, Inches(6.5), Inches(2.84), Inches(0.04), Inches(0.7), ACCENT_GREEN)
add_text_box(slide, Inches(4.2), Inches(3.55), Inches(2.5), Inches(0.35),
             "feature/IS-1234-signup", font_size=10, color=ACCENT_GREEN, font_name="Consolas")

add_accent_bar(slide, Inches(7.0), Inches(2.84), Inches(0.04), Inches(0.7), ACCENT_GREEN)
add_accent_bar(slide, Inches(7.0), Inches(3.5), Inches(2.5), Inches(0.04), ACCENT_GREEN)
add_accent_bar(slide, Inches(9.5), Inches(2.84), Inches(0.04), Inches(0.7), ACCENT_GREEN)
add_text_box(slide, Inches(7.2), Inches(3.55), Inches(2.5), Inches(0.35),
             "feature/IS-1235-verify", font_size=10, color=ACCENT_GREEN, font_name="Consolas")

# Labels
add_text_box(slide, Inches(2.7), Inches(2.1), Inches(0.5), Inches(0.35),
             "\u2193", font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(10.2), Inches(2.1), Inches(0.5), Inches(0.35),
             "\u2191", font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(1.6), Inches(1.5), Inches(0.3),
             "", font_size=10, color=MID_GRAY)
add_text_box(slide, Inches(10.6), Inches(1.65), Inches(1.2), Inches(0.3),
             "final PR", font_size=10, color=MID_GRAY)
add_text_box(slide, Inches(4.0), Inches(4.0), Inches(2.5), Inches(0.3),
             "\u2191 PR into integration branch", font_size=9, color=MID_GRAY)
add_text_box(slide, Inches(7.0), Inches(4.0), Inches(2.5), Inches(0.3),
             "\u2191 PR into integration branch", font_size=9, color=MID_GRAY)

# -- Rules (left) --
add_card(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.7))
add_accent_bar(slide, Inches(0.8), Inches(4.8), Inches(0.08), Inches(2.7), ACCENT_RED)

add_text_box(slide, Inches(1.2), Inches(4.9), Inches(4.8), Inches(0.4),
             "RULES", font_size=14, color=ACCENT_RED, bold=True)

ms_rules = [
    ("\u2713  integration/ prefix triggers this pattern", "CI detects integration/ and enforces multi-story rules"),
    ("\u2713  Story branches PR into integration branch", "feature/IS-xxxx \u2192 integration/IS-xxx (NOT into main)"),
    ("\u2713  Each story gets its own PR + review", "Traceability preserved at story level"),
    ("\u2713  Integration branch PRs into main/develop", "Final PR lists all included user stories"),
    ("\u2713  Integration branch is Standard Protected", "PR required, but branch is deleted after merge"),
    ("\u2713  Story branches deleted after merge", "Clean up as each story merges into integration"),
]
for i, (rule, detail) in enumerate(ms_rules):
    y = Inches(5.35) + i * Inches(0.35)
    add_text_box(slide, Inches(1.2), y, Inches(4.8), Inches(0.3),
                 rule, font_size=11, color=LIGHT_GRAY)

# -- When to use + why (right) --
add_card(slide, Inches(6.8), Inches(4.8), Inches(5.5), Inches(2.7))
add_accent_bar(slide, Inches(6.8), Inches(4.8), Inches(0.08), Inches(2.7), ACCENT_GREEN)

add_text_box(slide, Inches(7.2), Inches(4.9), Inches(4.8), Inches(0.4),
             "WHEN TO USE", font_size=14, color=ACCENT_GREEN, bold=True)

when_b = [
    "Feature decomposes into 2 or more user stories",
    "Multiple developers working on the same feature",
    "Stories need to be integrated before release",
    "Feature needs to be tested as a whole before merging to main",
]
for i, item in enumerate(when_b):
    add_text_box(slide, Inches(7.5), Inches(5.35) + i * Inches(0.38), Inches(4.5), Inches(0.35),
                 f"\u2713  {item}", font_size=12, color=LIGHT_GRAY)

add_text_box(slide, Inches(7.2), Inches(6.95), Inches(4.8), Inches(0.4),
             "WHY NOT JUST MERGE STORIES DIRECTLY?", font_size=11, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(7.2), Inches(7.25), Inches(4.8), Inches(0.3),
             "Stories may depend on each other. The feature\nbranch lets you integrate and test them together\nbefore they hit main.", font_size=10, color=LIGHT_GRAY)

add_slide_number(slide, 11, TOTAL_SLIDES)

# ============================================================
# SLIDE 12 — Traceability Chain
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "End-to-End Traceability Chain", font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Every artifact links back to the originating requirement", font_size=14, color=MID_GRAY)

chain = [
    ("Epic", "IS-100\nPlanning only\nNo branch", ACCENT_PURPLE),
    ("Feature", "IS-200\nintegration/ branch\n(if multi-story)", ACCENT_BLUE),
    ("User Story", "IS-1234\nfeature/ branch\n(default level)", ACCENT_GREEN),
    ("Branch", "feature/IS-1234-\nadd-profile", ACCENT_GREEN),
    ("Pull Request", "feat(user): IS-1234\nAdd user profile", ACCENT_ORANGE),
    ("Release", "release/v2.4.0\nIncludes: IS-1234", ACCENT_RED),
]

for i, (label, detail, color) in enumerate(chain):
    x = Inches(0.4) + i * Inches(2.1)
    y = Inches(1.8)

    add_card(slide, x, y, Inches(1.9), Inches(2.8))
    add_accent_bar(slide, x, y, Inches(1.9), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.2), Inches(1.7), Inches(0.8),
                 label, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.1), Inches(1.7), Inches(1.2),
                 detail, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER,
                 font_name="Consolas")

    if i < len(chain) - 1:
        add_text_box(slide, x + Inches(1.7), y + Inches(1.1), Inches(0.6), Inches(0.5),
                     "\u2192", font_size=24, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Audit query box
add_card(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(2.0))
add_text_box(slide, Inches(1.2), Inches(5.3), Inches(4), Inches(0.4),
             "AUDIT QUERIES THIS ENABLES", font_size=13, color=ACCENT_BLUE, bold=True)

queries = [
    ("\u2713  \"What code changed for IS-1234?\"  \u2192  Search branches & PRs by ticket ID",),
    ("\u2713  \"Who approved the change?\"  \u2192  PR approval history with reviewer names",),
    ("\u2713  \"Which release included this fix?\"  \u2192  Release branch contains merged PR",),
    ("\u2713  \"What was the business justification?\"  \u2192  PR links to user story with acceptance criteria",),
]

for i, (q,) in enumerate(queries):
    add_text_box(slide, Inches(1.5), Inches(5.7) + i * Inches(0.38), Inches(10), Inches(0.35),
                 q, font_size=13, color=LIGHT_GRAY)

add_slide_number(slide, 12, TOTAL_SLIDES)

# ============================================================
# SLIDE 13 — Commit Automation Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Commit Traceability Automation", font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Developers focus on code — tooling handles traceability", font_size=14, color=MID_GRAY)

# -- What's Required vs Optional --
add_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.3))
add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(0.08), Inches(2.3), ACCENT_RED)

add_text_box(slide, Inches(1.2), Inches(1.6), Inches(4.8), Inches(0.4),
             "REQUIRED (enforced by CI)", font_size=14, color=ACCENT_RED, bold=True)
req_items = [
    "Ticket ID in branch name",
    "Ticket ID in PR title",
    "Conventional commit format",
    "PR linked to ticket in description",
]
for i, item in enumerate(req_items):
    add_text_box(slide, Inches(1.5), Inches(2.1) + i * Inches(0.38), Inches(4.5), Inches(0.35),
                 f"\u2713  {item}", font_size=13, color=LIGHT_GRAY)

add_card(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(2.3))
add_accent_bar(slide, Inches(6.8), Inches(1.5), Inches(0.08), Inches(2.3), ACCENT_GREEN)

add_text_box(slide, Inches(7.2), Inches(1.6), Inches(4.8), Inches(0.4),
             "OPTIONAL (handled by automation)", font_size=14, color=ACCENT_GREEN, bold=True)
opt_items = [
    "Ticket ID in commit message body/footer",
    "Refs: footer in each commit",
    "Cross-referencing related tickets",
    "Linking commits to deployment artifacts",
]
for i, item in enumerate(opt_items):
    add_text_box(slide, Inches(7.5), Inches(2.1) + i * Inches(0.38), Inches(4.5), Inches(0.35),
                 f"\u2713  {item}", font_size=13, color=LIGHT_GRAY)

# -- Automation Pipeline --
add_text_box(slide, Inches(0.8), Inches(4.1), Inches(6), Inches(0.4),
             "AUTOMATION PIPELINE", font_size=13, color=ACCENT_PURPLE, bold=True)

pipeline_steps = [
    ("prepare-commit-msg\nGit Hook", "Runs locally on each commit.\nExtracts ticket ID from branch\nname and appends\nRefs: <ticket-id> to the\ncommit message footer.",
     ACCENT_PURPLE, "LOCAL"),
    ("Branch Name\nValidator", "CI step runs on PR open.\nRejects PRs from branches\nthat don't match the pattern:\n<type>/<TICKET-ID>-<desc>",
     ACCENT_BLUE, "CI"),
    ("PR Title\nValidator", "CI step checks PR title\nmatches conventional format\nwith ticket ID.\nBlocks merge if invalid.",
     ACCENT_ORANGE, "CI"),
    ("Merge Commit\nTagger", "On squash-merge, the merge\ncommit automatically includes\nthe PR number and ticket ID\nin the final commit message.",
     ACCENT_GREEN, "CI"),
]

for i, (title, desc, accent, where) in enumerate(pipeline_steps):
    x = Inches(0.5) + i * Inches(3.2)
    y = Inches(4.5)
    add_card(slide, x, y, Inches(2.9), Inches(2.8))
    add_accent_bar(slide, x, y, Inches(2.9), Inches(0.06), accent)

    where_badge = add_card(slide, x + Inches(1.9), y + Inches(0.15), Inches(0.8), Inches(0.35), accent)
    add_text_box(slide, x + Inches(1.9), y + Inches(0.17), Inches(0.8), Inches(0.3),
                 where, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(1.7), Inches(0.6),
                 title, font_size=14, color=accent, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.85), Inches(2.5), Inches(1.8),
                 desc, font_size=11, color=LIGHT_GRAY)

    if i < len(pipeline_steps) - 1:
        add_text_box(slide, x + Inches(2.7), y + Inches(1.0), Inches(0.7), Inches(0.5),
                     "\u2192", font_size=22, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 13, TOTAL_SLIDES)

# ============================================================
# SLIDE 14 — Quick Reference Card
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Quick Reference", font_size=32, color=WHITE, bold=True)

# Left column: Branch
add_card(slide, Inches(0.8), Inches(1.3), Inches(5.7), Inches(5.7))
add_accent_bar(slide, Inches(0.8), Inches(1.3), Inches(5.7), Inches(0.06), ACCENT_BLUE)
add_text_box(slide, Inches(1.2), Inches(1.5), Inches(5), Inches(0.5),
             "BRANCH NAMING", font_size=16, color=ACCENT_BLUE, bold=True)

branch_refs = [
    "feature/IS-xxxx-short-desc",
    "feature/pi3-s2/IS-xxxx-short-desc",
    "integration/IS-xxxx-short-desc",
    "bugfix/IS-xxxx-short-desc",
    "hotfix/IS-xxxx-short-desc",
    "chore/IS-xxxx-short-desc",
    "release/vX.Y.Z",
]
for i, ref in enumerate(branch_refs):
    add_text_box(slide, Inches(1.2), Inches(2.1) + i * Inches(0.45), Inches(5), Inches(0.4),
                 ref, font_size=14, color=WHITE, font_name="Consolas")

add_accent_bar(slide, Inches(1.2), Inches(4.5), Inches(4.8), Inches(0.02), CARD_BG)
add_text_box(slide, Inches(1.2), Inches(4.6), Inches(5), Inches(0.5),
             "PR TITLE", font_size=16, color=ACCENT_GREEN, bold=True)

pr_refs = [
    "feat(scope): IS-xxxx Description",
    "fix(scope): IS-xxxx Description",
    "hotfix(scope): IS-xxxx Description",
    "chore(scope): IS-xxxx Description",
]
for i, ref in enumerate(pr_refs):
    add_text_box(slide, Inches(1.2), Inches(5.2) + i * Inches(0.45), Inches(5), Inches(0.4),
                 ref, font_size=14, color=WHITE, font_name="Consolas")

# Right column: Do / Don't
add_card(slide, Inches(6.8), Inches(1.3), Inches(5.7), Inches(2.6))
add_accent_bar(slide, Inches(6.8), Inches(1.3), Inches(5.7), Inches(0.06), ACCENT_GREEN)
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5), Inches(0.5),
             "DO", font_size=16, color=ACCENT_GREEN, bold=True)

dos = [
    "\u2713  Include ticket ID in branch name",
    "\u2713  One ticket per branch / PR",
    "\u2713  Squash-merge feature branches",
    "\u2713  Delete branch after merge",
]
for i, d in enumerate(dos):
    add_text_box(slide, Inches(7.2), Inches(2.1) + i * Inches(0.4), Inches(5), Inches(0.4),
                 d, font_size=13, color=LIGHT_GRAY)

add_card(slide, Inches(6.8), Inches(4.2), Inches(5.7), Inches(2.8))
add_accent_bar(slide, Inches(6.8), Inches(4.2), Inches(5.7), Inches(0.06), ACCENT_RED)
add_text_box(slide, Inches(7.2), Inches(4.4), Inches(5), Inches(0.5),
             "DON'T", font_size=16, color=ACCENT_RED, bold=True)

donts = [
    "\u2717  Push directly to main or develop",
    "\u2717  Use vague branch names (my-fix, test-2)",
    "\u2717  Merge without PR approval",
    "\u2717  Leave stale branches after merge",
    "\u2717  Skip ticket ID in branch or PR title",
]
for i, d in enumerate(donts):
    add_text_box(slide, Inches(7.2), Inches(5.0) + i * Inches(0.4), Inches(5), Inches(0.4),
                 d, font_size=13, color=LIGHT_GRAY)

add_slide_number(slide, 14, TOTAL_SLIDES)

# ============================================================
# SLIDE 15 — Enforcement & Tooling
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Enforcement & Tooling", font_size=32, color=WHITE, bold=True)

tools = [
    ("GitHub Branch Protection", "Enforce PR reviews, status checks,\nand prevent force pushes on\nmain and develop.", ACCENT_BLUE,
     "Settings \u2192 Branches \u2192 Add Rule"),
    ("Branch Name Lint", "CI job validates branch name\nmatches regex pattern.\nBlocks PR if non-compliant.", ACCENT_GREEN,
     "^(feature|integration|bugfix|hotfix|chore|release)/(pi\\d+-s\\d+/)?"),
    ("PR Title Lint", "GitHub Action or CI step validates\nPR title follows conventional\ncommit format.", ACCENT_ORANGE,
     "commitlint or semantic-pr-title"),
    ("Commit Message Hook", "pre-commit hook or CI check\nenforces conventional commits.\nRejects non-compliant messages.", ACCENT_PURPLE,
     "commitlint + husky (local hooks)"),
]

for i, (title, desc, accent, tool_hint) in enumerate(tools):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.1)
    y = Inches(1.3) + row * Inches(2.7)
    add_card(slide, x, y, Inches(5.7), Inches(2.3))
    add_accent_bar(slide, x, y, Inches(0.08), Inches(2.3), accent)
    add_text_box(slide, x + Inches(0.4), y + Inches(0.25), Inches(5), Inches(0.5),
                 title, font_size=18, color=accent, bold=True)
    add_text_box(slide, x + Inches(0.4), y + Inches(0.8), Inches(5), Inches(1.0),
                 desc, font_size=14, color=LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.4), y + Inches(1.7), Inches(5), Inches(0.4),
                 tool_hint, font_size=12, color=MID_GRAY, font_name="Consolas")

# Bottom note
add_card(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5))
add_text_box(slide, Inches(1.2), Inches(6.82), Inches(10.5), Inches(0.4),
             "Automate enforcement wherever possible — convention without enforcement is just a suggestion.",
             font_size=14, color=ACCENT_ORANGE, bold=True)

add_slide_number(slide, 15, TOTAL_SLIDES)

# ============================================================
# Save
# ============================================================
output = "/Users/cuztomizer/Documents/GitHub/tacticl-core/docs/SDLC-Branch-Traceability-Standards.pptx"
prs.save(output)
print(f"Saved: {output}")
